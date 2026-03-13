#!/usr/bin/env python3
"""
ole_embedder.py - Embed CDXML files as editable ChemDraw OLE objects in PPTX/DOCX.

Takes one or more CDXML files and creates a PowerPoint or Word document with
each structure embedded as a double-clickable, editable ChemDraw OLE object.

Uses binary OLE construction with CDXML BoundingBox sizing for correct
display dimensions. Requires ChemDraw COM for CDX conversion and EMF preview.

Usage:
    # Single structure
    python ole_embedder.py scheme.cdxml --pptx -o report.pptx
    python ole_embedder.py scheme.cdxml --docx -o report.docx

    # Multiple structures (one per slide / one per paragraph)
    python ole_embedder.py s1.cdxml s2.cdxml s3.cdxml --pptx -o report.pptx

    # With margin adjustment (points)
    python ole_embedder.py scheme.cdxml --pptx --margin 2.0 -o report.pptx

Requirements:
    - ChemDraw 16+ (COM automation for CDXML -> CDX + EMF)
    - python-pptx (PPTX scaffolding)
    - python-docx (DOCX scaffolding)
    - lxml
"""

import argparse
import io
import json
import os
import struct
import sys
import tempfile
import xml.etree.ElementTree as ET
import zipfile


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# ChemDraw OLE CLSID (little-endian mixed-endian UUID encoding)
# UUID: 41BA6D21-A02E-11CE-8FD9-0020AFD1F20C
CHEMDRAW_CLSID = bytes([
    0x21, 0x6D, 0xBA, 0x41,
    0x2E, 0xA0,
    0xCE, 0x11,
    0x8F, 0xD9,
    0x00, 0x20, 0xAF, 0xD1, 0xF2, 0x0C,
])

SECTOR_SIZE = 512
MINI_SECTOR_SIZE = 64
MINI_STREAM_CUTOFF = 4096
ENDOFCHAIN = 0xFFFFFFFE
FREESECT = 0xFFFFFFFF
NOSTREAM = 0xFFFFFFFF
FATSECT = 0xFFFFFFFD

# Constant bytes from a known-good ChemDraw OLE object (via Office COM)
COMPOBJ_BYTES = bytes.fromhex(
    "0100feff030a0000ffffffff"
    "216dba412ea0ce118fd90020afd1f20c"
    "140000004353204368656d447261772044726177696e67"
    "001c0000004368656d4472617720496e7465726368616e676520466f726d6174"
    "00160000004368656d447261772e446f63756d656e742e362e30"
    "00f439b271000000000000000000000000"
)  # 126 bytes

OLE_BYTES = bytes.fromhex(
    "0100000200000000000000000000000000000000"
)  # 20 bytes

OLEPRES000_BYTES = bytes.fromhex(
    "ffffffff030000000400000001000000"
    "ffffffff02000000000000000000000000000000"
    "000000004e414e4900000000"
)  # 48 bytes


# ---------------------------------------------------------------------------
# ChemDraw COM — batch conversion
# ---------------------------------------------------------------------------

def batch_convert(cdxml_paths):
    """
    Open ChemDraw once and convert all CDXML files to CDX + EMF.

    Returns list of dicts: [{
        'path': str, 'name': str,
        'cdx_data': bytes, 'emf_data': bytes,
        'width_emu': int, 'height_emu': int,
    }, ...]
    """
    import win32com.client as win32

    results = []
    app = None
    try:
        app = win32.Dispatch("ChemDraw.Application")
        app.Visible = False

        for cdxml_path in cdxml_paths:
            cdxml_abs = os.path.abspath(cdxml_path)
            name = os.path.splitext(os.path.basename(cdxml_path))[0]

            tmp_cdx = tempfile.mktemp(suffix=".cdx")
            tmp_emf = tempfile.mktemp(suffix=".emf")

            doc = None
            try:
                doc = app.Documents.Open(cdxml_abs)

                doc.SaveAs(os.path.abspath(tmp_cdx))
                with open(tmp_cdx, "rb") as f:
                    cdx_data = f.read()

                doc.SaveAs(os.path.abspath(tmp_emf))
                with open(tmp_emf, "rb") as f:
                    emf_data = f.read()

                results.append({
                    'path': cdxml_path,
                    'name': name,
                    'cdx_data': cdx_data,
                    'emf_data': emf_data,
                })
            finally:
                try:
                    if doc:
                        doc.Close(False)
                except Exception:
                    pass
                for tmp in (tmp_cdx, tmp_emf):
                    if os.path.exists(tmp):
                        os.unlink(tmp)

    finally:
        try:
            if app:
                app.Quit()
        except Exception:
            pass

    return results


# ---------------------------------------------------------------------------
# CDXML dimension parsing
# ---------------------------------------------------------------------------

def get_cdxml_content_size(cdxml_path, margin_pt=0.0, scale=1.02):
    """
    Compute OLE display dimensions from CDXML content BoundingBox.

    The CDXML root BoundingBox is the tight content bounding box in points.
    A scale factor (default 1.02 = +2%) compensates for the slight difference
    between the CDXML BoundingBox and ChemDraw's actual OLE extent.
    Convert to EMU: 1 pt = 12700 EMU.

    Returns (width_emu, height_emu).
    """
    tree = ET.parse(cdxml_path)
    root = tree.getroot()

    w_pt = h_pt = None

    # Try root BoundingBox first (tightest content box)
    bb = root.get('BoundingBox')
    if bb:
        parts = bb.split()
        if len(parts) == 4:
            x1, y1, x2, y2 = map(float, parts)
            w_pt = (x2 - x1) + 2 * margin_pt
            h_pt = (y2 - y1) + 2 * margin_pt

    # Fallback: compute from all child BoundingBox attributes
    if w_pt is None:
        min_x, min_y = float("inf"), float("inf")
        max_x, max_y = float("-inf"), float("-inf")
        for elem in root.iter():
            child_bb = elem.get("BoundingBox")
            if child_bb and elem.tag not in ("page", "CDXML"):
                parts = child_bb.split()
                if len(parts) == 4:
                    try:
                        cx1, cy1, cx2, cy2 = map(float, parts)
                        min_x = min(min_x, cx1)
                        min_y = min(min_y, cy1)
                        max_x = max(max_x, cx2)
                        max_y = max(max_y, cy2)
                    except ValueError:
                        continue

        if min_x < float("inf"):
            w_pt = (max_x - min_x) + 2 * margin_pt
            h_pt = (max_y - min_y) + 2 * margin_pt

    # Last resort: use page dimensions (scale not applied)
    if w_pt is None:
        for page in root.iter('page'):
            pw = page.get('Width')
            ph = page.get('Height')
            if pw and ph:
                return int(float(pw) * 12700), int(float(ph) * 12700)
        return 5080000, 2540000  # ~4" x 2" fallback

    # Apply scale factor
    w_pt *= scale
    h_pt *= scale

    return int(w_pt * 12700), int(h_pt * 12700)


# ---------------------------------------------------------------------------
# CFB (Compound File Binary) builder
# ---------------------------------------------------------------------------

def _make_dir_entry(name, entry_type, color, left_id, right_id, child_id,
                    clsid, start_sector, size):
    """Build a 128-byte CFB directory entry."""
    entry = bytearray(128)

    name_utf16 = name.encode("utf-16-le") + b"\x00\x00"
    name_len = len(name_utf16)
    entry[0:name_len] = name_utf16

    struct.pack_into("<H", entry, 64, name_len)
    entry[66] = entry_type  # 2=stream, 5=root
    entry[67] = color       # 0=red, 1=black

    struct.pack_into("<I", entry, 68, left_id)
    struct.pack_into("<I", entry, 72, right_id)
    struct.pack_into("<I", entry, 76, child_id)

    if clsid:
        entry[80:96] = clsid

    struct.pack_into("<I", entry, 116, start_sector)
    struct.pack_into("<I", entry, 120, size & 0xFFFFFFFF)

    return bytes(entry)


def build_ole_compound_file(cdx_data):
    """
    Build a CFB file matching the known-good layout from Office COM.

    Layout:
        Header (512 bytes)
        Sector 0-1: Free
        Sector 2:   Directory sector 1 (entries 0-3)
        Sector 3:   FAT sector
        Sector 4:   Directory sector 2 (entry 4 + 3 empty)
        Sector 5:   Mini-FAT sector
        Sector 6:   Mini-stream container (256 bytes)
        Sectors 7+: CONTENTS data (CDX binary)
    """
    cdx_sectors = (len(cdx_data) + SECTOR_SIZE - 1) // SECTOR_SIZE
    cdx_start = 7
    total_sectors = cdx_start + cdx_sectors

    if total_sectors > 128:
        raise ValueError(f"CDX data too large ({len(cdx_data)} bytes) for single-FAT CFB")

    # --- FAT ---
    fat = [FREESECT] * 128
    fat[0] = FREESECT
    fat[1] = FREESECT
    fat[2] = 4              # dir sector 1 -> sector 4
    fat[3] = FATSECT
    fat[4] = ENDOFCHAIN     # dir sector 2 end
    fat[5] = ENDOFCHAIN     # mini-FAT
    fat[6] = ENDOFCHAIN     # mini-stream container

    for i in range(cdx_sectors):
        sec = cdx_start + i
        fat[sec] = sec + 1 if i < cdx_sectors - 1 else ENDOFCHAIN

    # --- Mini-FAT ---
    mini_fat = [FREESECT] * 128
    mini_fat[0] = ENDOFCHAIN   # OlePres000
    mini_fat[1] = ENDOFCHAIN   # Ole
    mini_fat[2] = 3            # CompObj -> 3
    mini_fat[3] = ENDOFCHAIN   # CompObj end

    # --- Mini-stream container (256 bytes) ---
    mini_stream = bytearray(256)
    mini_stream[0:48] = OLEPRES000_BYTES      # mini-sector 0
    mini_stream[64:84] = OLE_BYTES            # mini-sector 1
    mini_stream[128:254] = COMPOBJ_BYTES      # mini-sectors 2-3

    # --- Directory entries ---
    e0 = _make_dir_entry("Root Entry", 5, 0,
                         NOSTREAM, NOSTREAM, 1,
                         CHEMDRAW_CLSID, 6, 256)
    e1 = _make_dir_entry("\x01CompObj", 2, 1,
                         3, 2, NOSTREAM,
                         None, 2, 126)
    e2 = _make_dir_entry("CONTENTS", 2, 1,
                         NOSTREAM, 4, NOSTREAM,
                         None, cdx_start, len(cdx_data))
    e3 = _make_dir_entry("\x01Ole", 2, 1,
                         NOSTREAM, NOSTREAM, NOSTREAM,
                         None, 1, 20)
    e4 = _make_dir_entry("\x02OlePres000", 2, 1,
                         NOSTREAM, NOSTREAM, NOSTREAM,
                         None, 0, 48)

    dir_sector_1 = e0 + e1 + e2 + e3
    dir_sector_2 = e4 + (b"\x00" * 128 * 3)

    # --- CFB header ---
    header = bytearray(SECTOR_SIZE)
    header[0:8] = b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1"
    struct.pack_into("<H", header, 24, 0x003E)
    struct.pack_into("<H", header, 26, 0x0003)
    struct.pack_into("<H", header, 28, 0xFFFE)
    struct.pack_into("<H", header, 30, 9)
    struct.pack_into("<H", header, 32, 6)
    struct.pack_into("<I", header, 40, 0)
    struct.pack_into("<I", header, 44, 1)
    struct.pack_into("<I", header, 48, 2)
    struct.pack_into("<I", header, 52, 0)
    struct.pack_into("<I", header, 56, MINI_STREAM_CUTOFF)
    struct.pack_into("<I", header, 60, 5)
    struct.pack_into("<I", header, 64, 1)
    struct.pack_into("<I", header, 68, ENDOFCHAIN)
    struct.pack_into("<I", header, 72, 0)
    struct.pack_into("<I", header, 76, 3)
    for i in range(1, 109):
        struct.pack_into("<I", header, 76 + i * 4, FREESECT)

    # --- Assemble ---
    out = io.BytesIO()
    out.write(bytes(header))
    out.write(b"\x00" * SECTOR_SIZE)       # sector 0
    out.write(b"\x00" * SECTOR_SIZE)       # sector 1
    out.write(dir_sector_1)                # sector 2
    out.write(b"".join(struct.pack("<I", x) for x in fat))  # sector 3
    out.write(dir_sector_2)                # sector 4
    out.write(b"".join(struct.pack("<I", x) for x in mini_fat))  # sector 5
    ms_padded = bytes(mini_stream) + b"\x00" * (SECTOR_SIZE - len(mini_stream))
    out.write(ms_padded)                   # sector 6
    cdx_padded = cdx_data + b"\x00" * (cdx_sectors * SECTOR_SIZE - len(cdx_data))
    out.write(cdx_padded)                  # sectors 7+

    return out.getvalue()


# ---------------------------------------------------------------------------
# OOXML helpers (shared)
# ---------------------------------------------------------------------------

def _ensure_content_types(ct_xml):
    """Ensure .bin and .emf content types are declared."""
    from lxml import etree

    root = etree.fromstring(ct_xml)
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"

    existing = set()
    for elem in root.findall(f"{{{ct_ns}}}Default"):
        existing.add(elem.get("Extension", "").lower())

    if "bin" not in existing:
        etree.SubElement(root, f"{{{ct_ns}}}Default", attrib={
            "Extension": "bin",
            "ContentType": "application/vnd.openxmlformats-officedocument.oleObject",
        })
    if "emf" not in existing:
        etree.SubElement(root, f"{{{ct_ns}}}Default", attrib={
            "Extension": "emf",
            "ContentType": "image/x-emf",
        })

    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


# ---------------------------------------------------------------------------
# PPTX builder — one OLE object per slide
# ---------------------------------------------------------------------------

def build_pptx(items, output_path):
    """
    Create a PPTX with one editable ChemDraw OLE object per slide.

    items: list of dicts with 'ole_data', 'emf_data', 'width_emu', 'height_emu', 'name'
    """
    from pptx import Presentation

    prs = Presentation()

    # Find blank layout
    blank_layout = None
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[6]

    # Create one slide per item
    for _ in items:
        prs.slides.add_slide(blank_layout)

    tmp_pptx = tempfile.mktemp(suffix=".pptx")
    prs.save(tmp_pptx)

    try:
        _inject_pptx(tmp_pptx, output_path, items)
    finally:
        if os.path.exists(tmp_pptx):
            os.unlink(tmp_pptx)


def _inject_pptx(source_pptx, output_path, items):
    """Inject OLE objects into each slide of a PPTX."""
    from lxml import etree

    with zipfile.ZipFile(source_pptx, "r") as zin:
        with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zout:
            # Copy existing entries, modifying as needed
            for entry in zin.namelist():
                data = zin.read(entry)

                # Modify slide XML for each slide
                for idx, item in enumerate(items):
                    slide_num = idx + 1
                    slide_path = f"ppt/slides/slide{slide_num}.xml"
                    rels_path = f"ppt/slides/_rels/slide{slide_num}.xml.rels"
                    ole_idx = idx + 1

                    if entry == slide_path:
                        data = _build_slide_xml(data, item, ole_idx)
                    elif entry == rels_path:
                        data = _add_slide_rels(data, ole_idx)

                if entry == "[Content_Types].xml":
                    data = _ensure_content_types(data)

                zout.writestr(entry, data)

            # Add OLE + EMF files
            for idx, item in enumerate(items):
                ole_idx = idx + 1
                zout.writestr(
                    f"ppt/embeddings/oleObject{ole_idx}.bin",
                    item['ole_data'])
                zout.writestr(
                    f"ppt/media/olePreview{ole_idx}.emf",
                    item['emf_data'])


def _build_slide_xml(slide_xml, item, ole_idx):
    """Add OLE mc:AlternateContent to a slide."""
    from lxml import etree

    root = etree.fromstring(slide_xml)

    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    mc_ns = "http://schemas.openxmlformats.org/markup-compatibility/2006"
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    v_ns = "urn:schemas-microsoft-com:vml"

    sp_tree = root.find(f".//{{{p_ns}}}spTree")
    if sp_tree is None:
        raise RuntimeError("Could not find spTree in slide XML")

    w = item['width_emu']
    h = item['height_emu']

    # Center on slide (standard slide: 9144000 x 6858000 EMU)
    left = max(0, (9144000 - w) // 2)
    top = max(0, (6858000 - h) // 2)

    ole_rel = f"rIdOle{ole_idx}"
    img_rel = f"rIdOleImg{ole_idx}"
    base_id = 100 + (ole_idx - 1) * 10

    alt_xml = f"""<mc:AlternateContent
        xmlns:mc="{mc_ns}" xmlns:p="{p_ns}"
        xmlns:a="{a_ns}" xmlns:r="{r_ns}" xmlns:v="{v_ns}">
      <mc:Choice Requires="v">
        <p:graphicFrame>
          <p:nvGraphicFramePr>
            <p:cNvPr id="{base_id}" name="Object {ole_idx}"/>
            <p:cNvGraphicFramePr>
              <a:graphicFrameLocks noChangeAspect="1"/>
            </p:cNvGraphicFramePr>
            <p:nvPr/>
          </p:nvGraphicFramePr>
          <p:xfrm>
            <a:off x="{left}" y="{top}"/>
            <a:ext cx="{w}" cy="{h}"/>
          </p:xfrm>
          <a:graphic>
            <a:graphicData uri="http://schemas.openxmlformats.org/presentationml/2006/ole">
              <p:oleObj name="CS ChemDraw Drawing"
                        r:id="{ole_rel}"
                        imgW="{w}" imgH="{h}"
                        progId="ChemDraw.Document.6.0">
                <p:embed/>
              </p:oleObj>
            </a:graphicData>
          </a:graphic>
        </p:graphicFrame>
      </mc:Choice>
      <mc:Fallback>
        <p:graphicFrame>
          <p:nvGraphicFramePr>
            <p:cNvPr id="{base_id + 1}" name="Object {ole_idx}"/>
            <p:cNvGraphicFramePr>
              <a:graphicFrameLocks noChangeAspect="1"/>
            </p:cNvGraphicFramePr>
            <p:nvPr/>
          </p:nvGraphicFramePr>
          <p:xfrm>
            <a:off x="{left}" y="{top}"/>
            <a:ext cx="{w}" cy="{h}"/>
          </p:xfrm>
          <a:graphic>
            <a:graphicData uri="http://schemas.openxmlformats.org/presentationml/2006/ole">
              <p:oleObj name="CS ChemDraw Drawing"
                        r:id="{ole_rel}"
                        imgW="{w}" imgH="{h}"
                        progId="ChemDraw.Document.6.0">
                <p:embed/>
                <p:pic>
                  <p:nvPicPr>
                    <p:cNvPr id="{base_id + 2}" name="Preview {ole_idx}"/>
                    <p:cNvPicPr/>
                    <p:nvPr/>
                  </p:nvPicPr>
                  <p:blipFill>
                    <a:blip r:embed="{img_rel}"/>
                    <a:stretch><a:fillRect/></a:stretch>
                  </p:blipFill>
                  <p:spPr>
                    <a:xfrm>
                      <a:off x="{left}" y="{top}"/>
                      <a:ext cx="{w}" cy="{h}"/>
                    </a:xfrm>
                    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
                  </p:spPr>
                </p:pic>
              </p:oleObj>
            </a:graphicData>
          </a:graphic>
        </p:graphicFrame>
      </mc:Fallback>
    </mc:AlternateContent>"""

    sp_tree.append(etree.fromstring(alt_xml))
    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def _add_slide_rels(rels_xml, ole_idx):
    """Add OLE + image relationships to a slide's rels."""
    from lxml import etree

    root = etree.fromstring(rels_xml)

    etree.SubElement(root, "Relationship", attrib={
        "Id": f"rIdOle{ole_idx}",
        "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject",
        "Target": f"../embeddings/oleObject{ole_idx}.bin",
    })
    etree.SubElement(root, "Relationship", attrib={
        "Id": f"rIdOleImg{ole_idx}",
        "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
        "Target": f"../media/olePreview{ole_idx}.emf",
    })

    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


# ---------------------------------------------------------------------------
# DOCX builder — one OLE object per paragraph
# ---------------------------------------------------------------------------

def build_docx(items, output_path):
    """
    Create a DOCX with editable ChemDraw OLE objects, one per paragraph.

    items: list of dicts with 'ole_data', 'emf_data', 'width_emu', 'height_emu', 'name'
    """
    from docx import Document

    doc = Document()
    doc.add_paragraph("")  # placeholder
    tmp_docx = tempfile.mktemp(suffix=".docx")
    doc.save(tmp_docx)

    try:
        _inject_docx(tmp_docx, output_path, items)
    finally:
        if os.path.exists(tmp_docx):
            os.unlink(tmp_docx)


def _inject_docx(source_docx, output_path, items):
    """Inject OLE objects into a DOCX."""
    from lxml import etree

    with zipfile.ZipFile(source_docx, "r") as zin:
        with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for entry in zin.namelist():
                data = zin.read(entry)

                if entry == "word/document.xml":
                    data = _build_docx_xml(data, items)
                elif entry == "word/_rels/document.xml.rels":
                    data = _add_docx_rels(data, len(items))
                elif entry == "[Content_Types].xml":
                    data = _ensure_content_types(data)

                zout.writestr(entry, data)

            for idx, item in enumerate(items):
                ole_idx = idx + 1
                zout.writestr(
                    f"word/embeddings/oleObject{ole_idx}.bin",
                    item['ole_data'])
                zout.writestr(
                    f"word/media/olePreview{ole_idx}.emf",
                    item['emf_data'])


def _build_docx_xml(doc_xml, items):
    """Build document.xml with OLE objects."""
    from lxml import etree

    root = etree.fromstring(doc_xml)
    w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

    body = root.find(f"{{{w_ns}}}body")
    if body is None:
        raise RuntimeError("Could not find body in document XML")

    sect_pr = body.find(f"{{{w_ns}}}sectPr")

    for idx, item in enumerate(items):
        ole_idx = idx + 1
        w_emu = item['width_emu']
        h_emu = item['height_emu']

        w_pt = w_emu / 914400 * 72
        h_pt = h_emu / 914400 * 72
        w_twips = int(w_pt * 20)
        h_twips = int(h_pt * 20)

        shape_id = f"_x0000_s{1026 + idx}"
        object_id = f"_{1728379061 + idx}"

        para_xml = f"""<w:p xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
                            xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
                            xmlns:o="urn:schemas-microsoft-com:office:office"
                            xmlns:v="urn:schemas-microsoft-com:vml"
                            xmlns:w10="urn:schemas-microsoft-com:office:word">
            <w:r>
                <w:object w:dxaOrig="{w_twips}" w:dyaOrig="{h_twips}">
                    <v:shape id="{shape_id}"
                             type="#_x0000_t75"
                             style="width:{w_pt:.1f}pt;height:{h_pt:.1f}pt"
                             o:ole="">
                        <v:imagedata r:id="rIdOleImg{ole_idx}" o:title=""/>
                    </v:shape>
                    <o:OLEObject Type="Embed"
                                 ProgID="ChemDraw.Document.6.0"
                                 ShapeID="{shape_id}"
                                 DrawAspect="Content"
                                 ObjectID="{object_id}"
                                 r:id="rIdOle{ole_idx}"/>
                </w:object>
            </w:r>
        </w:p>"""

        ole_elem = etree.fromstring(para_xml)

        if sect_pr is not None:
            sect_pr.addprevious(ole_elem)
        else:
            body.append(ole_elem)

    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def _add_docx_rels(rels_xml, count):
    """Add OLE + image relationships for all objects."""
    from lxml import etree

    root = etree.fromstring(rels_xml)

    for ole_idx in range(1, count + 1):
        etree.SubElement(root, "Relationship", attrib={
            "Id": f"rIdOle{ole_idx}",
            "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject",
            "Target": f"embeddings/oleObject{ole_idx}.bin",
        })
        etree.SubElement(root, "Relationship", attrib={
            "Id": f"rIdOleImg{ole_idx}",
            "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
            "Target": f"media/olePreview{ole_idx}.emf",
        })

    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main(argv=None) -> int:
    parser = argparse.ArgumentParser(
        description="Embed CDXML files as editable ChemDraw OLE objects in PPTX or DOCX."
    )
    parser.add_argument("inputs", nargs="+", help="Input CDXML file(s)")
    parser.add_argument("-o", "--output", required=True, help="Output file path (.pptx or .docx)")
    parser.add_argument("--margin", type=float, default=0.0,
                        help="Margin in points around content (default: 0.0)")

    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--pptx", action="store_true", help="Create PowerPoint .pptx")
    group.add_argument("--docx", action="store_true", help="Create Word .docx")

    parser.add_argument("--json", action="store_true",
                        help="Output result as JSON to stdout")

    args = parser.parse_args(argv)

    # When --json, redirect status prints to stderr
    _log = print
    if args.json:
        def _log(*a, **kw):
            kw.setdefault("file", sys.stderr)
            print(*a, **kw)

    # Validate inputs
    for path in args.inputs:
        if not os.path.isfile(path):
            print(f"Error: file not found: {path}", file=sys.stderr)
            return 1
        if not path.lower().endswith(".cdxml"):
            print(f"Warning: {path} is not a .cdxml file", file=sys.stderr)

    ext = ".pptx" if args.pptx else ".docx"
    if not args.output.lower().endswith(ext):
        print(f"Warning: output file doesn't end with {ext}", file=sys.stderr)

    n = len(args.inputs)
    fmt = "PPTX" if args.pptx else "DOCX"
    _log(f"Embedding {n} structure{'s' if n > 1 else ''} into {fmt}...")

    # Step 1: Batch convert CDXML -> CDX + EMF
    _log(f"[1/3] Converting {n} CDXML file{'s' if n > 1 else ''} via ChemDraw COM...")
    converted = batch_convert(args.inputs)
    for item in converted:
        _log(f"       {item['name']}: CDX={len(item['cdx_data']):,}B, EMF={len(item['emf_data']):,}B")

    # Step 2: Compute dimensions + build OLE compound files
    _log(f"[2/3] Computing dimensions and building OLE files...")
    items = []
    for item in converted:
        w_emu, h_emu = get_cdxml_content_size(item['path'], margin_pt=args.margin)
        ole_data = build_ole_compound_file(item['cdx_data'])

        w_pt = w_emu / 12700
        h_pt = h_emu / 12700
        _log(f"       {item['name']}: {w_pt:.1f} x {h_pt:.1f} pt, OLE={len(ole_data):,}B")

        items.append({
            'ole_data': ole_data,
            'emf_data': item['emf_data'],
            'width_emu': w_emu,
            'height_emu': h_emu,
            'name': item['name'],
        })

    # Step 3: Build output document
    _log(f"[3/3] Building {fmt}...")
    if args.pptx:
        build_pptx(items, args.output)
    else:
        build_docx(items, args.output)

    if args.json:
        result = {
            "input_files": [os.path.abspath(p) for p in args.inputs],
            "output": os.path.abspath(args.output),
            "format": fmt.lower(),
            "num_objects_embedded": n,
        }
        print(json.dumps(result, indent=2))
    else:
        print(f"\nDone! {args.output}")
        print(f"  {n} editable ChemDraw OLE object{'s' if n > 1 else ''}")
        if args.pptx and n > 1:
            print(f"  {n} slides (one per structure)")

    return 0


if __name__ == "__main__":
    sys.exit(main())
