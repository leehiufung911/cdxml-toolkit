#!/usr/bin/env python3
"""
doc_from_template.py — Fill a PowerPoint or Word template with text and ChemDraw OLE structures.

Two-pass approach:
  Pass 1: python-pptx/python-docx replaces text placeholders (preserving formatting)
  Pass 2: XML-level injection replaces CDXML placeholders with editable OLE objects

Usage:
    python doc_from_template.py --template template.pptx --manifest manifest.json -o output.pptx
    python doc_from_template.py --template template.docx --manifest manifest.json -o output.docx
    python doc_from_template.py --create-test-template    # creates templates/reaction_summary.pptx

Requirements:
    - ChemDraw 16+ (COM automation for CDXML -> CDX + EMF) — only needed for cdxml slots
    - python-pptx
    - python-docx
    - lxml
"""

import argparse
import json
import os
import re
import shutil
import sys
import tempfile
import zipfile

from .ole_embedder import (
    batch_convert,
    get_cdxml_content_size,
    build_ole_compound_file,
    _ensure_content_types,
)


# ---------------------------------------------------------------------------
# XML namespaces
# ---------------------------------------------------------------------------

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
V_NS = "urn:schemas-microsoft-com:vml"
W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
O_NS = "urn:schemas-microsoft-com:office:office"
RELS_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
OLEOBJ_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject"
IMAGE_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"


# ---------------------------------------------------------------------------
# Manifest loading
# ---------------------------------------------------------------------------

def load_manifest(manifest_path):
    """Load JSON manifest. Resolve CDXML paths relative to manifest directory.

    Returns (text_slots, cdxml_slots, warnings).
    """
    with open(manifest_path) as f:
        data = json.load(f)

    base_dir = os.path.dirname(os.path.abspath(manifest_path))
    text_slots = {}    # placeholder -> value
    cdxml_slots = []   # [{"placeholder": ..., "file": abs_path}, ...]
    warnings = []

    for slot in data.get("slots", []):
        ph = slot["placeholder"]
        stype = slot.get("type", "text")

        if stype == "text":
            text_slots[ph] = slot["value"]
        elif stype == "cdxml":
            fpath = slot["file"]
            if not os.path.isabs(fpath):
                fpath = os.path.join(base_dir, fpath)
            fpath = os.path.abspath(fpath)
            if not os.path.isfile(fpath):
                warnings.append(f"CDXML file not found: {fpath}")
            cdxml_slots.append({"placeholder": ph, "file": fpath})
        else:
            warnings.append(f"Unknown slot type '{stype}' for {ph}")

    return text_slots, cdxml_slots, warnings


# ---------------------------------------------------------------------------
# OLE preparation (ChemDraw COM batch conversion)
# ---------------------------------------------------------------------------

def prepare_ole_items(cdxml_slots, margin_pt=0.0):
    """Convert unique CDXML files to OLE data via ChemDraw COM.

    Returns dict: abs_path -> {"ole_data", "emf_data", "width_emu", "height_emu"}
    """
    unique_files = list(dict.fromkeys(
        s["file"] for s in cdxml_slots if os.path.isfile(s["file"])
    ))
    if not unique_files:
        return {}

    converted = batch_convert(unique_files)

    items = {}
    for conv in converted:
        path = os.path.abspath(conv["path"])
        w_emu, h_emu = get_cdxml_content_size(conv["path"], margin_pt=margin_pt)
        ole_data = build_ole_compound_file(conv["cdx_data"])
        items[path] = {
            "ole_data": ole_data,
            "emf_data": conv["emf_data"],
            "width_emu": w_emu,
            "height_emu": h_emu,
        }

    return items


# ---------------------------------------------------------------------------
# Pass 1: Text replacement (python-pptx / python-docx)
# ---------------------------------------------------------------------------

def _replace_in_paragraph(paragraph, text_slots):
    """Replace {{PLACEHOLDER}} patterns in a paragraph's runs.

    Joins all run texts, performs replacements, puts result in first run.
    Preserves the first run's formatting. Returns set of filled placeholder names.
    """
    runs = paragraph.runs
    if not runs:
        return set()

    full_text = "".join(r.text or "" for r in runs)
    filled = set()
    new_text = full_text

    for placeholder, value in text_slots.items():
        if placeholder in new_text:
            new_text = new_text.replace(placeholder, value)
            filled.add(placeholder)

    if filled:
        runs[0].text = new_text
        for r in runs[1:]:
            r.text = ""

    return filled


def pass1_pptx(template_path, text_slots, temp_path):
    """Replace text placeholders in PPTX template. Save to temp_path.

    Returns set of placeholder names that were filled.
    """
    from pptx import Presentation

    prs = Presentation(template_path)
    filled = set()

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    filled.update(_replace_in_paragraph(para, text_slots))

    prs.save(temp_path)
    return filled


def pass1_docx(template_path, text_slots, temp_path):
    """Replace text placeholders in DOCX template. Save to temp_path.

    Returns set of placeholder names that were filled.
    """
    from docx import Document

    doc = Document(template_path)
    filled = set()

    for para in doc.paragraphs:
        filled.update(_replace_in_paragraph(para, text_slots))

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    filled.update(_replace_in_paragraph(para, text_slots))

    doc.save(temp_path)
    return filled


# ---------------------------------------------------------------------------
# Pass 2 helpers: relationship + content type management
# ---------------------------------------------------------------------------

def _rels_path_for(entry):
    """Compute the .rels file path for a ZIP entry (forward-slash paths)."""
    idx = entry.rfind("/")
    if idx < 0:
        return f"_rels/{entry}.rels"
    return f"{entry[:idx]}/_rels/{entry[idx + 1:]}.rels"


def _add_ole_rels(rels_xml, ole_idx, target_prefix):
    """Add OLE + image relationship entries to a rels XML document.

    target_prefix: '../' for PPTX slides, '' for DOCX document.
    """
    from lxml import etree

    root = etree.fromstring(rels_xml)

    etree.SubElement(root, "Relationship", attrib={
        "Id": f"rIdOle{ole_idx}",
        "Type": OLEOBJ_TYPE,
        "Target": f"{target_prefix}embeddings/oleObject{ole_idx}.bin",
    })
    etree.SubElement(root, "Relationship", attrib={
        "Id": f"rIdOleImg{ole_idx}",
        "Type": IMAGE_TYPE,
        "Target": f"{target_prefix}media/olePreview{ole_idx}.emf",
    })

    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


# ---------------------------------------------------------------------------
# Pass 2 — PPTX: replace CDXML placeholder shapes with OLE objects
# ---------------------------------------------------------------------------

def _make_pptx_ole_xml(ole_idx, x, y, w, h):
    """Build mc:AlternateContent XML for an OLE object in a PPTX slide."""
    from lxml import etree

    bid = 10000 + (ole_idx - 1) * 10
    orel = f"rIdOle{ole_idx}"
    irel = f"rIdOleImg{ole_idx}"

    xml_str = f"""<mc:AlternateContent
        xmlns:mc="{MC_NS}" xmlns:p="{P_NS}"
        xmlns:a="{A_NS}" xmlns:r="{R_NS}" xmlns:v="{V_NS}">
      <mc:Choice Requires="v">
        <p:graphicFrame>
          <p:nvGraphicFramePr>
            <p:cNvPr id="{bid}" name="ChemDraw {ole_idx}"/>
            <p:cNvGraphicFramePr>
              <a:graphicFrameLocks noChangeAspect="1"/>
            </p:cNvGraphicFramePr>
            <p:nvPr/>
          </p:nvGraphicFramePr>
          <p:xfrm>
            <a:off x="{x}" y="{y}"/>
            <a:ext cx="{w}" cy="{h}"/>
          </p:xfrm>
          <a:graphic>
            <a:graphicData uri="http://schemas.openxmlformats.org/presentationml/2006/ole">
              <p:oleObj name="CS ChemDraw Drawing" r:id="{orel}"
                        imgW="{w}" imgH="{h}" progId="ChemDraw.Document.6.0">
                <p:embed/>
              </p:oleObj>
            </a:graphicData>
          </a:graphic>
        </p:graphicFrame>
      </mc:Choice>
      <mc:Fallback>
        <p:graphicFrame>
          <p:nvGraphicFramePr>
            <p:cNvPr id="{bid + 1}" name="ChemDraw {ole_idx}"/>
            <p:cNvGraphicFramePr>
              <a:graphicFrameLocks noChangeAspect="1"/>
            </p:cNvGraphicFramePr>
            <p:nvPr/>
          </p:nvGraphicFramePr>
          <p:xfrm>
            <a:off x="{x}" y="{y}"/>
            <a:ext cx="{w}" cy="{h}"/>
          </p:xfrm>
          <a:graphic>
            <a:graphicData uri="http://schemas.openxmlformats.org/presentationml/2006/ole">
              <p:oleObj name="CS ChemDraw Drawing" r:id="{orel}"
                        imgW="{w}" imgH="{h}" progId="ChemDraw.Document.6.0">
                <p:embed/>
                <p:pic>
                  <p:nvPicPr>
                    <p:cNvPr id="{bid + 2}" name="Preview {ole_idx}"/>
                    <p:cNvPicPr/><p:nvPr/>
                  </p:nvPicPr>
                  <p:blipFill>
                    <a:blip r:embed="{irel}"/>
                    <a:stretch><a:fillRect/></a:stretch>
                  </p:blipFill>
                  <p:spPr>
                    <a:xfrm>
                      <a:off x="{x}" y="{y}"/>
                      <a:ext cx="{w}" cy="{h}"/>
                    </a:xfrm>
                    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
                  </p:spPr>
                </p:pic>
              </p:oleObj>
            </a:graphicData>
          </a:graphic>
        </p:graphicFrame>
      </mc:Fallback>
    </mc:AlternateContent>"""

    return etree.fromstring(xml_str)


def pass2_pptx(input_path, output_path, cdxml_slots, ole_items):
    """Replace CDXML placeholder text boxes with OLE objects in PPTX.

    Returns (ole_count, filled_placeholders_set).
    """
    from lxml import etree

    ph_to_file = {s["placeholder"]: s["file"] for s in cdxml_slots}
    ole_counter = 0
    filled_ph = set()

    # Scan slides for placeholder shapes, build modifications
    # mods: slide_entry -> (modified_xml_bytes, [(ole_idx, item)])
    mods = {}

    with zipfile.ZipFile(input_path, "r") as zin:
        for entry in zin.namelist():
            if not re.match(r"ppt/slides/slide\d+\.xml$", entry):
                continue

            root = etree.fromstring(zin.read(entry))
            sp_tree = root.find(f".//{{{P_NS}}}spTree")
            if sp_tree is None:
                continue

            slide_oles = []

            for sp in list(sp_tree.findall(f"{{{P_NS}}}sp")):
                texts = [t.text for t in sp.iter(f"{{{A_NS}}}t") if t.text]
                full = "".join(texts).strip()

                for ph, fpath in ph_to_file.items():
                    if ph not in full or fpath not in ole_items:
                        continue

                    ole_counter += 1
                    item = ole_items[fpath]
                    filled_ph.add(ph)

                    # Get shape position from its transform
                    x = y = 0
                    xfrm = sp.find(f"{{{P_NS}}}spPr/{{{A_NS}}}xfrm")
                    if xfrm is None:
                        xfrm = sp.find(f".//{{{A_NS}}}xfrm")
                    if xfrm is not None:
                        off = xfrm.find(f"{{{A_NS}}}off")
                        if off is not None:
                            x = int(off.get("x", "0"))
                            y = int(off.get("y", "0"))

                    # Remove placeholder text box, add OLE graphic frame
                    sp_tree.remove(sp)
                    sp_tree.append(_make_pptx_ole_xml(
                        ole_counter, x, y,
                        item["width_emu"], item["height_emu"],
                    ))
                    slide_oles.append((ole_counter, item))
                    break  # one placeholder per shape

            if slide_oles:
                mods[entry] = (
                    etree.tostring(root, xml_declaration=True,
                                   encoding="UTF-8", standalone=True),
                    slide_oles,
                )

    # Build rels update map: rels_entry -> [(ole_idx, item)]
    rels_updates = {}
    for slide_entry, (_, oles) in mods.items():
        rp = _rels_path_for(slide_entry)
        rels_updates[rp] = oles

    # Write output ZIP
    with zipfile.ZipFile(input_path, "r") as zin:
        with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for entry in zin.namelist():
                data = zin.read(entry)

                # Swap in modified slide XML
                if entry in mods:
                    data = mods[entry][0]

                # Add OLE relationships to affected slide rels
                if entry in rels_updates:
                    for idx, _ in rels_updates[entry]:
                        data = _add_ole_rels(data, idx, "../")

                # Ensure .bin and .emf content types exist
                if entry == "[Content_Types].xml":
                    data = _ensure_content_types(data)

                zout.writestr(entry, data)

            # Write OLE + EMF binary files
            for _, (_, oles) in mods.items():
                for idx, item in oles:
                    zout.writestr(
                        f"ppt/embeddings/oleObject{idx}.bin", item["ole_data"])
                    zout.writestr(
                        f"ppt/media/olePreview{idx}.emf", item["emf_data"])

    return ole_counter, filled_ph


# ---------------------------------------------------------------------------
# Pass 2 — DOCX: replace CDXML placeholder paragraphs with OLE objects
# ---------------------------------------------------------------------------

def _make_docx_ole_para(ole_idx, w_emu, h_emu):
    """Build a DOCX paragraph containing a ChemDraw OLE object."""
    from lxml import etree

    w_pt = w_emu / 12700
    h_pt = h_emu / 12700
    w_twips = int(w_pt * 20)
    h_twips = int(h_pt * 20)
    shape_id = f"_x0000_s{1026 + ole_idx}"
    obj_id = f"_{1728379061 + ole_idx}"

    xml_str = f"""<w:p xmlns:w="{W_NS}"
                       xmlns:r="{R_NS}"
                       xmlns:o="{O_NS}"
                       xmlns:v="{V_NS}">
        <w:r>
            <w:object w:dxaOrig="{w_twips}" w:dyaOrig="{h_twips}">
                <v:shape id="{shape_id}" type="#_x0000_t75"
                         style="width:{w_pt:.1f}pt;height:{h_pt:.1f}pt"
                         o:ole="">
                    <v:imagedata r:id="rIdOleImg{ole_idx}" o:title=""/>
                </v:shape>
                <o:OLEObject Type="Embed"
                             ProgID="ChemDraw.Document.6.0"
                             ShapeID="{shape_id}"
                             DrawAspect="Content"
                             ObjectID="{obj_id}"
                             r:id="rIdOle{ole_idx}"/>
            </w:object>
        </w:r>
    </w:p>"""

    return etree.fromstring(xml_str)


def pass2_docx(input_path, output_path, cdxml_slots, ole_items):
    """Replace CDXML placeholder paragraphs with OLE objects in DOCX.

    Returns (ole_count, filled_placeholders_set).
    """
    from lxml import etree

    ph_to_file = {s["placeholder"]: s["file"] for s in cdxml_slots}
    ole_counter = 0
    filled_ph = set()
    oles = []  # [(ole_idx, item)]

    with zipfile.ZipFile(input_path, "r") as zin:
        doc_xml = zin.read("word/document.xml")
        root = etree.fromstring(doc_xml)
        body = root.find(f"{{{W_NS}}}body")
        if body is None:
            shutil.copy2(input_path, output_path)
            return 0, set()

        # Search body paragraphs for CDXML placeholders
        for p_elem in list(body.findall(f"{{{W_NS}}}p")):
            texts = [t.text for t in p_elem.iter(f"{{{W_NS}}}t") if t.text]
            full = "".join(texts).strip()

            for ph, fpath in ph_to_file.items():
                if ph not in full or fpath not in ole_items:
                    continue

                ole_counter += 1
                item = ole_items[fpath]
                filled_ph.add(ph)

                new_p = _make_docx_ole_para(
                    ole_counter, item["width_emu"], item["height_emu"])
                body.replace(p_elem, new_p)
                oles.append((ole_counter, item))
                break  # one placeholder per paragraph

    if not oles:
        shutil.copy2(input_path, output_path)
        return 0, filled_ph

    new_doc_xml = etree.tostring(
        root, xml_declaration=True, encoding="UTF-8", standalone=True)

    # Write output ZIP
    with zipfile.ZipFile(input_path, "r") as zin:
        with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for entry in zin.namelist():
                data = zin.read(entry)

                if entry == "word/document.xml":
                    data = new_doc_xml

                if entry == "word/_rels/document.xml.rels":
                    for idx, _ in oles:
                        data = _add_ole_rels(data, idx, "")

                if entry == "[Content_Types].xml":
                    data = _ensure_content_types(data)

                zout.writestr(entry, data)

            for idx, item in oles:
                zout.writestr(
                    f"word/embeddings/oleObject{idx}.bin", item["ole_data"])
                zout.writestr(
                    f"word/media/olePreview{idx}.emf", item["emf_data"])

    return ole_counter, filled_ph


# ---------------------------------------------------------------------------
# Test template creation
# ---------------------------------------------------------------------------

def create_test_template(output_dir="templates"):
    """Create a minimal 1-slide PPTX template with placeholder text boxes."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    os.makedirs(output_dir, exist_ok=True)

    prs = Presentation()

    # Find blank layout
    blank = None
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            blank = layout
            break
    if blank is None:
        blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)

    # Title text box
    tb = slide.shapes.add_textbox(Inches(1), Inches(0.4), Inches(8), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = "{{TITLE}}"
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle text box
    tb = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "{{SUBTITLE}}"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER

    # Scheme placeholder text box (positioned for a reaction scheme)
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
    p = tb.text_frame.paragraphs[0]
    p.text = "{{SCHEME}}"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.CENTER

    out_path = os.path.join(output_dir, "reaction_summary.pptx")
    prs.save(out_path)
    return out_path


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main(argv=None) -> int:
    parser = argparse.ArgumentParser(
        description="Fill a PowerPoint or Word template with text and "
                    "ChemDraw OLE structures from a JSON manifest."
    )
    parser.add_argument("--template", help="Template file (.pptx or .docx)")
    parser.add_argument("--manifest", help="JSON manifest file")
    parser.add_argument("-o", "--output", help="Output file path")
    parser.add_argument("--margin", type=float, default=0.0,
                        help="OLE margin in points (default: 0.0)")
    parser.add_argument("--json", action="store_true",
                        help="JSON output summary")
    parser.add_argument("--create-test-template", action="store_true",
                        help="Create templates/reaction_summary.pptx and exit")

    args = parser.parse_args(argv)

    # --create-test-template mode
    if args.create_test_template:
        out = create_test_template()
        if args.json:
            print(json.dumps({"template": out}))
        else:
            print(f"Created test template: {out}")
        return 0

    # Normal mode: validate required args
    if not args.template or not args.manifest or not args.output:
        parser.error("--template, --manifest, and -o are required")

    if not os.path.isfile(args.template):
        print(f"Error: template not found: {args.template}", file=sys.stderr)
        return 1
    if not os.path.isfile(args.manifest):
        print(f"Error: manifest not found: {args.manifest}", file=sys.stderr)
        return 1

    is_pptx = args.template.lower().endswith(".pptx")
    is_docx = args.template.lower().endswith(".docx")
    if not (is_pptx or is_docx):
        print("Error: template must be .pptx or .docx", file=sys.stderr)
        return 1

    # Load manifest
    text_slots, cdxml_slots, warnings = load_manifest(args.manifest)

    if not args.json:
        fmt = "PPTX" if is_pptx else "DOCX"
        print(f"Template: {args.template} ({fmt})")
        print(f"Manifest: {len(text_slots)} text slot(s), "
              f"{len(cdxml_slots)} CDXML slot(s)")

    # Pass 1: text replacement via python-pptx/python-docx
    ext = ".pptx" if is_pptx else ".docx"
    tmp = tempfile.mktemp(suffix=ext)

    try:
        if not args.json:
            print("[1/2] Replacing text placeholders...")

        if is_pptx:
            text_filled = pass1_pptx(args.template, text_slots, tmp)
        else:
            text_filled = pass1_docx(args.template, text_slots, tmp)

        # Warn about unfilled text slots
        for ph in text_slots:
            if ph not in text_filled:
                warnings.append(f"Text placeholder not found in template: {ph}")

        # Pass 2: OLE embedding
        ole_count = 0
        cdxml_filled = set()

        if cdxml_slots:
            if not args.json:
                print("[2/2] Converting CDXML and embedding OLE objects...")

            ole_items = prepare_ole_items(cdxml_slots, margin_pt=args.margin)

            # Warn about conversion failures
            for slot in cdxml_slots:
                if os.path.isfile(slot["file"]) and slot["file"] not in ole_items:
                    warnings.append(f"Failed to convert CDXML: {slot['file']}")

            if ole_items:
                if is_pptx:
                    ole_count, cdxml_filled = pass2_pptx(
                        tmp, args.output, cdxml_slots, ole_items)
                else:
                    ole_count, cdxml_filled = pass2_docx(
                        tmp, args.output, cdxml_slots, ole_items)
            else:
                shutil.copy2(tmp, args.output)
        else:
            if not args.json:
                print("[2/2] No CDXML slots — skipping OLE embedding.")
            shutil.copy2(tmp, args.output)

        # Warn about unfilled CDXML slots
        for slot in cdxml_slots:
            if slot["placeholder"] not in cdxml_filled:
                if slot["file"] in (ole_items if cdxml_slots else {}):
                    warnings.append(
                        f"CDXML placeholder not found in template: "
                        f"{slot['placeholder']}")

        total_filled = len(text_filled) + ole_count

        if args.json:
            result = {
                "template": args.template,
                "output": args.output,
                "slots_filled": total_filled,
                "ole_objects": ole_count,
                "warnings": warnings,
            }
            print(json.dumps(result, indent=2))
        else:
            print(f"\nDone! {args.output}")
            print(f"  Text slots filled: {len(text_filled)}")
            print(f"  OLE objects embedded: {ole_count}")
            if warnings:
                print(f"  Warnings:")
                for w in warnings:
                    print(f"    - {w}")

    finally:
        if os.path.exists(tmp):
            os.unlink(tmp)

    return 0


if __name__ == "__main__":
    sys.exit(main())
